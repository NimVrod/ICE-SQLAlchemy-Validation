"""Build ICE SQLAlchemy Migration Validation slide deck (.pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("ice-sqlalchemy-validation-slides.pptx")

ACCENT = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF9, 0xFA)
CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.55)
MARGIN_R = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
TITLE_TOP = Inches(0.4)
BODY_TOP = Inches(1.55)
BODY_H = Inches(5.55)
FOOTER_TOP = Inches(6.95)


def _fill_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    *,
    font_pt: int = 18,
    bold: bool = False,
    color: RGBColor = ACCENT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tf


def _add_footer(slide, page: int, total: int) -> None:
    inst = "Data Quality in Database Migrations"
    authors = "Hubert Dec · Patryk Dziki · Adrian Czapka · Dominik Dziedzic · Beniamin Sujka"
    w = CONTENT_W / 3
    for text, x, align in [
        (inst, MARGIN_L, PP_ALIGN.LEFT),
        (f"{page} / {total}", MARGIN_L + w, PP_ALIGN.CENTER),
        (authors, MARGIN_L + 2 * w, PP_ALIGN.RIGHT),
    ]:
        tf = _textbox(slide, x, FOOTER_TOP, w, Inches(0.35), font_pt=8, color=MUTED, align=align)
        tf.paragraphs[0].text = text


def _set_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    tf = _textbox(slide, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(0.65), font_pt=28, bold=True)
    tf.paragraphs[0].text = title
    if subtitle:
        tf2 = _textbox(
            slide,
            MARGIN_L,
            Inches(1.05),
            CONTENT_W,
            Inches(0.45),
            font_pt=16,
            color=MUTED,
        )
        tf2.paragraphs[0].text = subtitle
        tf2.paragraphs[0].font.italic = True


def _body_box(slide, top=BODY_TOP, height=BODY_H):
    return _textbox(slide, MARGIN_L, top, CONTENT_W, height, font_pt=16, color=ACCENT)


def _add_bullets(tf, items: list[str], *, level0_pt: int = 16) -> None:
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(level0_pt)
        p.font.color.rgb = ACCENT
        p.space_after = Pt(6)


def _content_slide(prs, title: str, subtitle: str | None, builder) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    _set_slide_title(slide, title, subtitle)
    builder(slide)
    return slide


def _add_table(slide, left, top, width, height, headers: list[str], rows: list[list[str]]):
    cols = len(headers)
    table_shape = slide.shapes.add_table(1 + len(rows), cols, left, top, width, height)
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = ACCENT
    return table


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    y = Inches(1.8)
    for text, size, bold, color in [
        ("DATA QUALITY IN DATABASE MIGRATIONS", 11, False, MUTED),
        ("ICE SQLAlchemy Migration Validation", 36, True, ACCENT),
        ("Post-migration data quality on existing rows", 20, False, MUTED),
    ]:
        tf = _textbox(
            slide,
            MARGIN_L,
            y,
            CONTENT_W,
            Inches(0.8),
            font_pt=size,
            bold=bold,
            color=color,
            align=PP_ALIGN.CENTER,
        )
        tf.paragraphs[0].text = text
        y += Inches(0.75 if bold else 0.55)
    tf = _textbox(
        slide,
        MARGIN_L,
        y + Inches(0.3),
        CONTENT_W,
        Inches(0.4),
        font_pt=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    tf.paragraphs[0].text = (
        "Hubert Dec · Patryk Dziki · Adrian Czapka · Dominik Dziedzic · Beniamin Sujka"
    )
    y += Inches(0.7)
    tf = _textbox(
        slide,
        MARGIN_L,
        y,
        CONTENT_W,
        Inches(1.2),
        font_pt=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    tf.paragraphs[0].text = (
        "Python library for auditing rows after schema migration — "
        "not at INSERT time, but on data already in the target database.\n\n"
        "Based on When Migration Changes Your Data (Sprint I report)"
    )

    # Slide 2
    def s2(slide):
        intro = _textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.7), font_pt=15)
        intro.paragraphs[0].text = (
            "Our Sprint I report studied real incidents (insurance, banking, travel). "
            "Pipelines stay green and row counts match, yet business logic breaks silently."
        )
        col_w = (CONTENT_W - Inches(0.3)) / 2
        left_tf = _body_box(slide, top=Inches(2.35), height=Inches(4.5))
        left_tf.paragraphs[0].text = "What the report shows"
        left_tf.paragraphs[0].font.bold = True
        left_tf.paragraphs[0].font.size = Pt(14)
        _add_bullets(
            left_tf,
            [
                "Row counts match; migration tools report success",
                "Columns stay non-null and syntactically valid",
                "Semantic drift: dates, categories, foreign keys, ownership",
                "Detection in production — tickets, audits, regulators",
                "Fixing later costs rollback, re-migration, or repair scripts",
            ],
        )
        right = slide.shapes.add_textbox(MARGIN_L + col_w + Inches(0.3), Inches(2.35), col_w, Inches(4.5))
        rtf = right.text_frame
        rtf.paragraphs[0].text = "Examples from our report"
        rtf.paragraphs[0].font.bold = True
        rtf.paragraphs[0].font.size = Pt(14)
        for line in [
            "[MEDIUM] Insurance dates → year 2040",
            "[HIGH] 87% of timestamps shifted",
            "[CRITICAL] TSB: wrong customer accounts",
            "[CRITICAL] TUI: category → wrong aircraft weight",
            "",
            "Systematic mapping or timezone errors that pass basic QA.",
        ]:
            p = rtf.add_paragraph()
            p.text = line
            p.font.size = Pt(13 if line.startswith("[") else 11)
            p.font.color.rgb = MUTED if not line.startswith("[") else ACCENT
            p.space_after = Pt(4)

    _content_slide(
        prs,
        "Why this library exists",
        "Migrations often succeed while data becomes wrong.",
        s2,
    )

    # Slide 3
    def s3(slide):
        intro = _textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.65), font_pt=15)
        intro.paragraphs[0].text = (
            "Teams assume ORM models and DB constraints are enough. After bulk load, "
            "millions of legacy rows may never be checked against the new rules."
        )
        _add_table(
            slide,
            MARGIN_L,
            Inches(2.2),
            CONTENT_W,
            Inches(2.2),
            ["Check", "Catches new bad INSERT", "Audits migrated rows"],
            [
                ["CheckConstraint on model", "Often", "Only if DB enforces + re-checked"],
                ["ORM @validates", "On flush", "No — not on bulk / raw SQL"],
                ["Row count vs source", "—", "Volume only"],
                ["MigrationValidator", "—", "Yes — scans existing data"],
            ],
        )
        foot = _textbox(slide, MARGIN_L, Inches(4.55), CONTENT_W, Inches(1.5), font_pt=14, color=MUTED)
        foot.paragraphs[0].text = (
            "Bulk import (pandas.to_sql, ETL, COPY, raw SQL) can leave violations in place. "
            "Teams often use lax staging tables and only find issues on explicit audit."
        )

    _content_slide(
        prs,
        "What usual checks miss",
        "Schema constraints protect new writes — not legacy rows already in the table.",
        s3,
    )

    # Slide 4
    def s4(slide):
        intro = _textbox(slide, MARGIN_L, Inches(1.15), CONTENT_W, Inches(0.9), font_pt=15)
        intro.paragraphs[0].text = (
            "Library in src/ validates data already in the database against SQLAlchemy 2.0 models. "
            "Same idea as unit tests — executed as set-based SQL over the full table."
        )
        col_w = (CONTENT_W - Inches(0.4)) / 3
        for i, (head, bullets) in enumerate(
            [
                (
                    "Input",
                    [
                        "Session against migrated DB",
                        "Mapped model (User, Policy, …)",
                        "Lax table OK; strict model",
                        "Mapped, Column, Annotated",
                    ],
                ),
                (
                    "Process",
                    [
                        "parser.py reads Annotated fields",
                        "Validators: null, bounds, NaN, unique",
                        "SQL SELECT with yield_per=5000",
                        "Duplicate pass via GROUP BY",
                        "Python only when SQL cannot (NaN)",
                    ],
                ),
                (
                    "Output",
                    [
                        "list[ValidationError]",
                        "errors_to_dataframe report",
                        "summarize_errors counts",
                        "Pytest plugin for CI",
                    ],
                ),
            ]
        ):
            tf = slide.shapes.add_textbox(
                MARGIN_L + i * (col_w + Inches(0.2)),
                Inches(2.15),
                col_w,
                Inches(4.6),
            ).text_frame
            tf.paragraphs[0].text = head
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(14)
            _add_bullets(tf, bullets, level0_pt=13)

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide4)
    _set_slide_title(slide4, "ICE SQLAlchemy Migration Validation", None)
    s4(slide4)

    # Slide 5
    def s5(slide):
        intro = _textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.55), font_pt=14)
        intro.paragraphs[0].text = (
            "Rules live on the column — same diff as the migration model, no YAML drift."
        )
        code = """from typing import Annotated
from sqlalchemy.orm import Mapped, mapped_column
from src.metadata import FloatMax, FloatMin, Required, Unique, UniqueComposite

class User(Base):
    __tablename__ = "users"
    email: Mapped[Annotated[str, Required(), Unique()]] = mapped_column(String(100))
    age: Mapped[Annotated[int, Required(), FloatMin(0), FloatMax(120)]] = mapped_column(Integer)
    country: Mapped[Annotated[str, Required(), UniqueComposite("country", "city")]] = mapped_column(String(40))
    city: Mapped[Annotated[str, Required()]] = mapped_column(String(40))"""
        box = slide.shapes.add_shape(
            1, MARGIN_L, Inches(2.0), CONTENT_W, Inches(2.85)  # MSO_SHAPE.RECTANGLE = 1
        )
        box.fill.solid()
        box.fill.fore_color.rgb = CODE_BG
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        tf = box.text_frame
        tf.paragraphs[0].text = code
        tf.paragraphs[0].font.name = "Consolas"
        tf.paragraphs[0].font.size = Pt(10)
        notes = _textbox(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(1.4), font_pt=12)
        _add_bullets(
            notes,
            [
                "Mapped[T | None] or Nullable() → NULL allowed",
                "Non-optional Mapped[T] implies Required() even if DB column is nullable",
                "UniqueComposite on one column validates tuple (country, city)",
            ],
            level0_pt=12,
        )

    _content_slide(
        prs,
        "One place to declare rules: the column",
        "No separate config file. No magic __table_args__ scanning.",
        s5,
    )

    # Slide 6
    def s6(slide):
        _add_table(
            slide,
            MARGIN_L,
            BODY_TOP,
            CONTENT_W,
            Inches(3.2),
            ["Marker", "Meaning", "Strategy"],
            [
                ["Required()", "Must not be NULL", "SQL: IS NULL"],
                ["Nullable()", "Allow NULL", "Skips required rule"],
                ["Unique()", "Column unique", "SQL GROUP BY + join"],
                ["UniqueComposite(a,b)", "Composite unique", "SQL duplicate scan"],
                ["FloatMin / FloatMax", "Numeric bounds", "SQL comparison"],
                ["NoNaN()", "No NaN when set", "SQL prefilter + Python isnan"],
            ],
        )
        foot = _textbox(slide, MARGIN_L, Inches(4.5), CONTENT_W, Inches(1.6), font_pt=13, color=MUTED)
        foot.paragraphs[0].text = (
            "Out of scope for v1: semantic transforms (e.g. category 3 → Boeing 737). "
            "Focus: NULLs, ranges, duplicates, NaN — failures inside successful migrations."
        )

    _content_slide(
        prs,
        "Validation markers",
        "Each marker maps to one validator; engine picks SQL or short Python follow-up.",
        s6,
    )

    # Slide 7
    def s7(slide):
        flow = _textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.6), font_pt=15)
        flow.paragraphs[0].text = (
            "Model (Annotated) → parser.py → SQL SELECT (stream) → ValidationError + pandas report"
        )
        tf = _body_box(slide, top=Inches(2.2), height=Inches(4.2))
        tf.paragraphs[0].text = "Field pass"
        tf.paragraphs[0].font.bold = True
        _add_bullets(
            tf,
            [
                "One validator per annotated column",
                "ValidationError includes primary key + column",
            ],
        )
        p = tf.add_paragraph()
        p.text = "Duplicate pass"
        p.font.bold = True
        p.font.size = Pt(14)
        p.space_before = Pt(12)
        _add_bullets(
            tf,
            [
                "Collects Unique() and UniqueComposite targets",
                "Single SQL pass per constraint group",
            ],
        )

    _content_slide(
        prs,
        "How MigrationValidator runs",
        "Field pass first, then duplicates. Errors stream; memory stays bounded.",
        s7,
    )

    # Slide 8
    def s8(slide):
        _add_table(
            slide,
            MARGIN_L,
            BODY_TOP,
            CONTENT_W,
            Inches(1.4),
            ["Validator", "SQL", "Python"],
            [
                ["Null / FloatMin / Max / Unique", "Full rule in SQL", "Confirm rule fired"],
                ["NoNaN", "Non-null prefilter", "math.isnan (SQLite-safe)"],
            ],
        )
        tf = _body_box(slide, top=Inches(2.65), height=Inches(3.8))
        tf.paragraphs[0].text = "Why it matters"
        tf.paragraphs[0].font.bold = True
        _add_bullets(
            tf,
            [
                "yield_per=5000 — not row-by-row ORM",
                "Lax staging table + strict model",
                "Same report for Faker, ETL, or legacy export",
                "Actionable errors: row id, column, rule kind",
            ],
            level0_pt=14,
        )

    _content_slide(
        prs,
        "Validate smartly: SQL first, Python when needed",
        "Push work to the database; Python only when SQL cannot express the rule.",
        s8,
    )

    # Slide 9
    def s9(slide):
        code = """from sqlalchemy.orm import Session
from src.engine import MigrationValidator
from src.pandas_ import errors_to_dataframe, summarize_errors

with Session(engine) as session:
    errors = MigrationValidator(session).validate(User)
    report = errors_to_dataframe(errors)
    summary = summarize_errors(report)"""
        box = slide.shapes.add_shape(1, MARGIN_L, BODY_TOP, CONTENT_W, Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = CODE_BG
        box.text_frame.paragraphs[0].text = code
        box.text_frame.paragraphs[0].font.name = "Consolas"
        box.text_frame.paragraphs[0].font.size = Pt(11)
        col_w = (CONTENT_W - Inches(0.3)) / 2
        for i, (head, bullets) in enumerate(
            [
                (
                    "Pytest",
                    [
                        "assert_table_valid(session, User)",
                        "pytest_validation_plugin in pyproject.toml",
                        "migration_validator fixture",
                    ],
                ),
                (
                    "Examples",
                    [
                        "examples/migration_validation_demo.py",
                        "examples/validate_dirty_database.py",
                        "examples/validation_report.ipynb",
                        "uv run pytest (29 tests)",
                    ],
                ),
            ]
        ):
            tf = slide.shapes.add_textbox(
                MARGIN_L + i * (col_w + Inches(0.3)),
                Inches(3.55),
                col_w,
                Inches(2.5),
            ).text_frame
            tf.paragraphs[0].text = head
            tf.paragraphs[0].font.bold = True
            _add_bullets(tf, bullets, level0_pt=13)

    _content_slide(
        prs,
        "Running validation",
        "Three lines in app code; optional CI gate via pytest.",
        s9,
    )

    # Slide 10
    def s10(slide):
        col_w = CONTENT_W * 0.55
        tf = slide.shapes.add_textbox(MARGIN_L, BODY_TOP, col_w, Inches(4.8)).text_frame
        tf.paragraphs[0].text = "Workflow"
        tf.paragraphs[0].font.bold = True
        steps = [
            "Generate ~1000 rows with Faker",
            "Inject defects: 50 NULL emails, 20 negative ages, 25 duplicates",
            "df.to_sql into SQLite (like ETL)",
            "MigrationValidator.validate(ValidatedUser)",
            "Print summary / notebook for demo",
        ]
        _add_bullets(tf, steps, level0_pt=14)
        box = slide.shapes.add_shape(
            1, MARGIN_L + col_w + Inches(0.25), BODY_TOP, CONTENT_W - col_w - Inches(0.25), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = CODE_BG
        box.text_frame.paragraphs[0].text = (
            "Found 122 validation error(s)\n\n"
            "null       email   52\n"
            "duplicate  email   50\n"
            "float_min  age     20"
        )
        box.text_frame.paragraphs[0].font.name = "Consolas"
        box.text_frame.paragraphs[0].font.size = Pt(12)
        foot = _textbox(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(1.2), font_pt=13, color=MUTED)
        foot.paragraphs[0].text = (
            "Data present and often non-null but wrong. Library produces evidence (row ids + rule) before sign-off."
        )

    _content_slide(
        prs,
        "Demo: dirty data like production migrations",
        "Local SQLite reproduces “looks fine in the dashboard” failures.",
        s10,
    )

    # Slide 11
    def s11(slide):
        col_w = (CONTENT_W - Inches(0.3)) / 2
        for i, (head, bullets) in enumerate(
            [
                (
                    "Run after",
                    [
                        "Bulk load / cutover on staging",
                        "Staging mirrors production volume",
                        "Before UAT and regulatory sign-off",
                        "After hotfix scripts on legacy data",
                    ],
                ),
                (
                    "Pair with report themes",
                    [
                        "Historical edge cases → custom markers",
                        "Code lists → FloatMin/Max or enums",
                        "Duplicates → Unique / UniqueComposite",
                        "Silent NULLs → Required() on nullable cols",
                    ],
                ),
            ]
        ):
            tf = slide.shapes.add_textbox(
                MARGIN_L + i * (col_w + Inches(0.3)),
                BODY_TOP,
                col_w,
                Inches(4.2),
            ).text_frame
            tf.paragraphs[0].text = head
            tf.paragraphs[0].font.bold = True
            _add_bullets(tf, bullets, level0_pt=14)
        goal = _textbox(slide, MARGIN_L, Inches(5.35), CONTENT_W, Inches(0.9), font_pt=15, bold=True)
        goal.paragraphs[0].text = (
            "Goal: detect in staging, not in production. "
            "Failing pytest beats a post-go-live incident report."
        )

    _content_slide(
        prs,
        "When to run this in a migration project",
        "Validation gate — same weight as schema diff and row-count checks.",
        s11,
    )

    # Slide 12
    def s12(slide):
        tf = _body_box(slide, top=Inches(1.7), height=Inches(4.5))
        for i, line in enumerate(
            [
                "Problem: migrations succeed while data is wrong (Sprint I report).",
                "Approach: Annotated markers on columns; one declaration.",
                "Engine: MigrationValidator + SQL streaming + Python for NaN.",
                "Delivery: pandas report, pytest plugin, Faker demo.",
                "Next: domain semantics (codes, FK) via same parser.",
            ],
            start=1,
        ):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"{i}. {line}"
            p.font.size = Pt(17)
            p.space_after = Pt(10)
        cmd = _textbox(
            slide,
            MARGIN_L,
            Inches(5.85),
            CONTENT_W,
            Inches(0.5),
            font_pt=13,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        cmd.paragraphs[0].text = (
            "uv run pytest  ·  uv run python examples/validate_dirty_database.py"
        )

    _content_slide(
        prs,
        "Summary",
        "Structural checks today; semantic rules tomorrow.",
        s12,
    )

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        _add_footer(slide, i - 1, total - 1)

    return prs


def main() -> None:
    prs = build()
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
